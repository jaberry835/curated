"""Text extraction service for different file types."""

import io
import json
import mimetypes
from typing import Dict, Any, Optional, Tuple
import chardet

try:
    from ..utils.logging import get_logger
except ImportError:
    try:
        from utils.logging import get_logger
    except ImportError:
        import logging
        def get_logger(name):
            return logging.getLogger(name)

logger = get_logger(__name__)


class TextExtractionService:
    """Service for extracting text content from various file types."""
    
    def __init__(self):
        """Initialize the text extraction service."""
        self.supported_types = {
            'text/plain': self._extract_text,
            'text/csv': self._extract_text,
            'text/markdown': self._extract_text,
            'text/html': self._extract_html,
            'application/json': self._extract_json,
            'application/xml': self._extract_xml,
            'text/xml': self._extract_xml,
            'application/pdf': self._extract_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._extract_xlsx,
            'application/msword': self._extract_doc,
            'application/vnd.ms-powerpoint': self._extract_ppt,
            'application/vnd.ms-excel': self._extract_xls
        }
        logger.info("Text extraction service initialized")
    
    def extract_text(self, file_content: bytes, content_type: str, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text content from file bytes.
        
        Args:
            file_content: The binary content of the file
            content_type: MIME type of the file
            file_name: Name of the file (used for fallback type detection)
            
        Returns:
            Tuple of (extracted_text, metadata)
        """
        try:
            # Normalize content type
            content_type = content_type.lower() if content_type else ''
            
            # Fallback: detect content type from filename if not provided or generic
            if not content_type or content_type in ['application/octet-stream', 'binary/octet-stream']:
                detected_type, _ = mimetypes.guess_type(file_name)
                if detected_type:
                    content_type = detected_type.lower()
                    logger.info(f"Detected content type from filename: {content_type}")
            
            # Get extraction function
            extract_func = self.supported_types.get(content_type, self._extract_fallback_text)
            
            # Extract text
            text, metadata = extract_func(file_content, file_name)
            
            # Add extraction metadata
            metadata.update({
                'content_type': content_type,
                'extracted_chars': len(text),
                'extraction_method': extract_func.__name__
            })
            
            logger.info(f"Extracted {len(text)} characters from {file_name} using {extract_func.__name__}")
            return text, metadata
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_name}: {str(e)}")
            return f"[Error extracting text from {file_name}: {str(e)}]", {'error': str(e)}
    
    def _extract_text(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from plain text files."""
        try:
            # Detect encoding
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            confidence = detected.get('confidence', 0)
            
            # Try detected encoding first, fallback to utf-8
            try:
                text = content.decode(encoding)
            except (UnicodeDecodeError, LookupError):
                text = content.decode('utf-8', errors='replace')
                encoding = 'utf-8'
                confidence = 0
            
            return text, {
                'encoding': encoding,
                'encoding_confidence': confidence,
                'lines': len(text.splitlines())
            }
            
        except Exception as e:
            return f"[Error reading text file: {str(e)}]", {'error': str(e)}
    
    def _extract_html(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from HTML files."""
        try:
            # Try to import BeautifulSoup for better HTML parsing
            try:
                from bs4 import BeautifulSoup
                
                # Detect encoding
                detected = chardet.detect(content)
                encoding = detected.get('encoding', 'utf-8')
                
                html_content = content.decode(encoding, errors='replace')
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Extract title
                title = soup.title.string if soup.title else file_name
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                
                # Get text
                text = soup.get_text(separator=' ', strip=True)
                
                return text, {
                    'title': title,
                    'encoding': encoding,
                    'has_title': bool(soup.title),
                    'extraction_library': 'beautifulsoup4'
                }
                
            except ImportError:
                # Fallback: simple text extraction
                html_content = content.decode('utf-8', errors='replace')
                # Basic HTML tag removal (very simple)
                import re
                text = re.sub(r'<[^>]+>', ' ', html_content)
                text = re.sub(r'\s+', ' ', text).strip()
                
                return text, {
                    'encoding': 'utf-8',
                    'extraction_library': 'regex_fallback'
                }
                
        except Exception as e:
            return f"[Error reading HTML file: {str(e)}]", {'error': str(e)}
    
    def _extract_json(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from JSON files."""
        try:
            text_content = content.decode('utf-8', errors='replace')
            data = json.loads(text_content)
            
            # Convert JSON to readable text
            text = json.dumps(data, indent=2, ensure_ascii=False)
            
            return text, {
                'json_valid': True,
                'json_type': type(data).__name__,
                'json_keys': list(data.keys()) if isinstance(data, dict) else None
            }
            
        except json.JSONDecodeError as e:
            # Return raw text if JSON is invalid
            text_content = content.decode('utf-8', errors='replace')
            return text_content, {
                'json_valid': False,
                'json_error': str(e)
            }
        except Exception as e:
            return f"[Error reading JSON file: {str(e)}]", {'error': str(e)}
    
    def _extract_xml(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from XML files."""
        try:
            import xml.etree.ElementTree as ET
            
            text_content = content.decode('utf-8', errors='replace')
            root = ET.fromstring(text_content)
            
            # Extract all text content
            text_parts = []
            for elem in root.iter():
                if elem.text:
                    text_parts.append(elem.text.strip())
                if elem.tail:
                    text_parts.append(elem.tail.strip())
            
            text = ' '.join(filter(None, text_parts))
            
            return text, {
                'xml_valid': True,
                'root_tag': root.tag,
                'element_count': len(list(root.iter()))
            }
            
        except ET.ParseError as e:
            # Return raw text if XML is invalid
            text_content = content.decode('utf-8', errors='replace')
            return text_content, {
                'xml_valid': False,
                'xml_error': str(e)
            }
        except Exception as e:
            return f"[Error reading XML file: {str(e)}]", {'error': str(e)}
    
    def _extract_pdf(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PDF files."""
        try:
            # Try to import PyPDF2 or pdfplumber
            try:
                import PyPDF2
                
                pdf_file = io.BytesIO(content)
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                
                text_parts = []
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_parts.append(page_text)
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num}: {str(e)}")
                
                text = '\n\n'.join(text_parts)
                
                return text, {
                    'page_count': len(pdf_reader.pages),
                    'extraction_library': 'PyPDF2',
                    'pages_with_text': len(text_parts)
                }
                
            except ImportError:
                return f"[PDF text extraction requires PyPDF2 library. File: {file_name}]", {
                    'error': 'PyPDF2 not available',
                    'file_type': 'PDF'
                }
                
        except Exception as e:
            return f"[Error reading PDF file: {str(e)}]", {'error': str(e)}
    
    def _extract_docx(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from DOCX files."""
        try:
            from docx import Document
            
            doc_file = io.BytesIO(content)
            doc = Document(doc_file)
            
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            text = '\n'.join(text_parts)
            
            return text, {
                'paragraph_count': len(doc.paragraphs),
                'extraction_library': 'python-docx'
            }
            
        except ImportError:
            return f"[DOCX text extraction requires python-docx library. File: {file_name}]", {
                'error': 'python-docx not available',
                'file_type': 'DOCX'
            }
        except Exception as e:
            return f"[Error reading DOCX file: {str(e)}]", {'error': str(e)}
    
    def _extract_xlsx(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from XLSX files."""
        try:
            import openpyxl
            
            excel_file = io.BytesIO(content)
            workbook = openpyxl.load_workbook(excel_file, data_only=True)
            
            text_parts = []
            sheet_info = {}
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_text = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        sheet_text.append('\t'.join(row_text))
                
                if sheet_text:
                    text_parts.append(f"Sheet: {sheet_name}\n" + '\n'.join(sheet_text))
                    sheet_info[sheet_name] = len(sheet_text)
            
            text = '\n\n'.join(text_parts)
            
            return text, {
                'sheet_count': len(workbook.sheetnames),
                'sheet_info': sheet_info,
                'extraction_library': 'openpyxl'
            }
            
        except ImportError:
            return f"[XLSX text extraction requires openpyxl library. File: {file_name}]", {
                'error': 'openpyxl not available',
                'file_type': 'XLSX'
            }
        except Exception as e:
            return f"[Error reading XLSX file: {str(e)}]", {'error': str(e)}
    
    def _extract_pptx(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            
            ppt_file = io.BytesIO(content)
            presentation = Presentation(ppt_file)
            
            text_parts = []
            slide_count = 0
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text_parts.append(f"Slide {slide_num + 1}:\n" + '\n'.join(slide_text))
                    slide_count += 1
            
            text = '\n\n'.join(text_parts)
            
            return text, {
                'slide_count': len(presentation.slides),
                'slides_with_text': slide_count,
                'extraction_library': 'python-pptx'
            }
            
        except ImportError:
            return f"[PPTX text extraction requires python-pptx library. File: {file_name}]", {
                'error': 'python-pptx not available',
                'file_type': 'PPTX'
            }
        except Exception as e:
            return f"[Error reading PPTX file: {str(e)}]", {'error': str(e)}
    
    def _extract_doc(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from legacy DOC files."""
        return f"[Legacy DOC format not supported. Please convert to DOCX. File: {file_name}]", {
            'error': 'Legacy DOC format not supported',
            'file_type': 'DOC'
        }
    
    def _extract_ppt(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from legacy PPT files."""
        return f"[Legacy PPT format not supported. Please convert to PPTX. File: {file_name}]", {
            'error': 'Legacy PPT format not supported',
            'file_type': 'PPT'
        }
    
    def _extract_xls(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Extract text from legacy XLS files."""
        return f"[Legacy XLS format not supported. Please convert to XLSX. File: {file_name}]", {
            'error': 'Legacy XLS format not supported',
            'file_type': 'XLS'
        }
    
    def _extract_fallback_text(self, content: bytes, file_name: str) -> Tuple[str, Dict[str, Any]]:
        """Fallback text extraction for unsupported file types."""
        try:
            # Try to decode as text
            detected = chardet.detect(content)
            encoding = detected.get('encoding', 'utf-8')
            
            if detected.get('confidence', 0) > 0.7:
                try:
                    text = content.decode(encoding, errors='replace')
                    return text, {
                        'extraction_method': 'fallback_text',
                        'encoding': encoding,
                        'confidence': detected.get('confidence', 0)
                    }
                except Exception:
                    pass
            
            # If text extraction fails, return a placeholder
            return f"[Binary file - text extraction not supported for this file type. File: {file_name}]", {
                'extraction_method': 'fallback_binary',
                'file_size': len(content)
            }
            
        except Exception as e:
            return f"[Error in fallback text extraction: {str(e)}]", {'error': str(e)}


# Global service instance
text_extraction_service = TextExtractionService()
